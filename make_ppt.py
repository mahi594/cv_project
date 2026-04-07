from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "NeuralVision Engine"
    subtitle.text = "Interactive Computer Vision Platform\n\nProject Presentation"
    
    # Slide 2: Problem Statement / Introduction
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Introduction"
    tf = body_shape.text_frame
    tf.text = "Learning Computer Vision algorithms is often abstract and mathematical."
    p = tf.add_paragraph()
    p.text = "Students and developers need a way to visualize algorithms dynamically."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "NeuralVision Engine bridges this gap by providing real-time, interactive feedback."
    p.level = 1
    
    # Slide 3: Key Features
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Key Features (Syllabus Driven)"
    tf = body_shape.text_frame
    tf.text = "The platform covers core CV syllabus topics:"
    
    topics = [
        "Color Space Representation (HSV, LAB, YCrCb)",
        "Fourier Analysis & High Pass Filtering",
        "Noise Reduction (Gaussian, Median, Bilateral)",
        "Edge Detection (Canny, LoG, DoG)",
        "Harris Corner Detection",
        "Feature Descriptors (ORB)"
    ]
    for topic in topics:
        p = tf.add_paragraph()
        p.text = topic
        p.level = 1

    # Slide 4: Technology Stack
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Technology Stack"
    tf = body_shape.text_frame
    
    tf.text = "Frontend"
    p = tf.add_paragraph()
    p.text = "React & Vite for high-performance rendering"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "TailwindCSS & Lucide-React for modern UI components"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Backend"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "FastAPI for robust, asynchronous API endpoints"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "OpenCV & NumPy for heavy computational CV algorithms"
    p.level = 1

    # Slide 5: How It Works
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "How It Works"
    tf = body_shape.text_frame
    tf.text = "1. User uploads an image via the web interface."
    p = tf.add_paragraph()
    p.text = "2. Selects a syllabus module and adjusts parameters."
    p = tf.add_paragraph()
    p.text = "3. Frontend sends image and parameters to FastAPI backend."
    p = tf.add_paragraph()
    p.text = "4. OpenCV processes the image and extracts intermediate algorithm states."
    p = tf.add_paragraph()
    p.text = "5. Visualizations are returned and displayed interactively on the UI."

    # Slide 6: Future Enhancements
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Future Enhancements"
    tf = body_shape.text_frame
    tf.text = "Add deep learning modules (e.g., Object Detection, Segmentation)."
    p = tf.add_paragraph()
    p.text = "Allow users to build custom CV execution pipelines."
    p = tf.add_paragraph()
    p.text = "Deploy to cloud (AWS/GCP) for global accessibility."

    # Slide 7: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You!"
    subtitle.text = "Questions?"

    ppt_path = "NeuralVision_Engine_Presentation.pptx"
    prs.save(ppt_path)
    print(f"Presentation saved successfully to {ppt_path}")

if __name__ == "__main__":
    create_presentation()
